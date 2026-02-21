from pathlib import Path
from typing import List
from collections import defaultdict
import pptx
from .parser_base import BaseParser

class PPTParser(BaseParser):
    """
    Parser for PPTX files
    """
    type = 'pptx'
    
    def _extract_content(self) -> str:
        """
        Extract content from PPTX file using python-pptx
        """
        try:
            prs = pptx.Presentation(self.file_path)
            full_text = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            
            return '\n'.join(full_text)
        except Exception as e:
            print(f"Error reading PPTX file {self.file_path}: {e}")
            return ""

    def _extract_metadata(self) -> defaultdict:
        metadata = defaultdict(str)
        try:
            prs = pptx.Presentation(self.file_path)
            core_props = prs.core_properties
            metadata['title'] = core_props.title
            metadata['author'] = core_props.author
            metadata['subject'] = core_props.subject
            metadata['keywords'] = core_props.keywords
            metadata['created'] = str(core_props.created)
            metadata['modified'] = str(core_props.modified)
        except Exception:
            pass
        return metadata

    def _check_format(self) -> bool:
        f_path: Path = Path(self.file_path)
        return f_path.exists() and f_path.suffix.lower() == '.pptx'
