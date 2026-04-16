import os
import pptx
from .TextChunkProcessor import TextChunkProcessor


class PPTParser(TextChunkProcessor):
    """
    PPT/PPTX 文件解析器
    支持 .ppt (旧版 PowerPoint 97-2003) 和 .pptx (新版)
    默认使用 SlideChunking（按幻灯片分块）
    """
    type = 'ppt'

    def __init__(self, file_path: str, **kwargs):
        super().__init__(file_path, **kwargs)
        self._slide_data = []  # 存储幻灯片数据供后续使用

    def _extract_content(self) -> str:
        """
        提取 PPT 内容，返回带幻灯片标记的文本格式：

        [SLIDE:1:标题1]
        内容1-1
        内容1-2
        [NOTES:]
        备注内容
        [SLIDE:2:标题2]
        内容2-1
        """
        ext = os.path.splitext(self.file_path)[1].lower()

        if ext == '.pptx':
            return self._extract_pptx()
        elif ext == '.ppt':
            return self._extract_ppt()
        else:
            return ""

    def _extract_pptx(self) -> str:
        """使用 python-pptx 提取 .pptx 文本"""
        try:
            prs = pptx.Presentation(self.file_path)
            self._slide_data = []
            full_text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                title = ""
                notes = ""

                # 提取标题（通常在占位符中）
                for shape in slide.shapes:
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                        if shape.placeholder_format.type in (1, 2, 3):  # TITLE, CENTER_TITLE, SUBTITLE
                            if shape.text.strip():
                                title = shape.text.strip()
                                break

                # 提取所有文本内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # 跳过已经作为标题的文本
                        if shape.text.strip() == title:
                            continue
                        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                            continue
                        slide_texts.append(shape.text.strip())

                # 提取备注
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text.strip():
                        notes = notes_frame.text.strip()

                # 构建幻灯片文本
                slide_parts = []

                # 幻灯片标记
                if title:
                    slide_parts.append(f"[SLIDE:{slide_num}:{title}]")
                else:
                    slide_parts.append(f"[SLIDE:{slide_num}]")

                # 内容
                if slide_texts:
                    slide_parts.append('\n'.join(slide_texts))

                # 备注
                if notes:
                    slide_parts.append(f"[NOTES:]\n{notes}")

                slide_text = '\n'.join(slide_parts)

                # 存储幻灯片数据
                self._slide_data.append({
                    'num': slide_num,
                    'title': title,
                    'content': '\n'.join(slide_texts),
                    'notes': notes
                })

                full_text_parts.append(slide_text)

            return '\n\n'.join(full_text_parts)
        except Exception as e:
            print(f"Error reading PPTX file {self.file_path}: {e}")
            return ""

    def _extract_ppt(self) -> str:
        """使用 pywin32 + PowerPoint COM 提取 .ppt 文本"""
        try:
            import win32com.client

            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            ppt_app.Visible = True

            try:
                abs_path = os.path.abspath(self.file_path)
                presentation = ppt_app.Presentations.Open(abs_path, ReadOnly=True, WithWindow=True)

                self._slide_data = []
                full_text_parts = []

                for slide_num, slide in enumerate(presentation.Slides, 1):
                    slide_texts = []
                    title = ""
                    notes = ""

                    for shape in slide.Shapes:
                        # 提取标题
                        try:
                            if shape.HasTextFrame:
                                if hasattr(shape, "PlaceholderFormat") and shape.PlaceholderFormat:
                                    try:
                                        if shape.PlaceholderFormat.Type in (1, 2, 3):  # TITLE types
                                            text = shape.TextFrame.TextRange.Text
                                            if text.strip():
                                                title = text.strip()
                                                continue
                                    except Exception:
                                        # Shape 不是占位符，跳过
                                        pass

                                # 提取文本内容
                                text = shape.TextFrame.TextRange.Text
                                if text.strip():
                                    slide_texts.append(text.strip())
                        except Exception:
                            # 某些形状可能无法访问 TextFrame，跳过
                            continue

                    # 提取备注
                    if slide.NotesPage.Shapes.Count > 0:
                        notes_shape = slide.NotesPage.Shapes.Placeholders(2)
                        if notes_shape:
                            notes = notes_shape.TextFrame.TextRange.Text.strip()

                    # 构建幻灯片文本
                    slide_parts = []
                    if title:
                        slide_parts.append(f"[SLIDE:{slide_num}:{title}]")
                    else:
                        slide_parts.append(f"[SLIDE:{slide_num}]")

                    if slide_texts:
                        slide_parts.append('\n'.join(slide_texts))

                    if notes:
                        slide_parts.append(f"[NOTES:]\n{notes}")

                    slide_text = '\n'.join(slide_parts)

                    self._slide_data.append({
                        'num': slide_num,
                        'title': title,
                        'content': '\n'.join(slide_texts),
                        'notes': notes
                    })

                    full_text_parts.append(slide_text)

                presentation.Close()
                return '\n\n'.join(full_text_parts)

            finally:
                ppt_app.Quit()

        except Exception as e:
            print(f"Error reading PPT file {self.file_path} with pywin32: {e}")
            return ""

    def get_slide_data(self):
        """获取所有幻灯片的结构化数据"""
        return self._slide_data
